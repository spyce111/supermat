import re
import json
import logging
from io import BytesIO
import zipfile
import spacy
# import nltk
from tqdm import tqdm
from .constants import *
import uuid
import traceback
import string
import difflib
import os
# import numpy as np
# from transformers.pipelines import AggregationStrategy
from adobe.pdfservices.operation.auth.credentials import Credentials
from adobe.pdfservices.operation.exception.exceptions import ServiceApiException, ServiceUsageException, SdkException
from adobe.pdfservices.operation.pdfops.options.extractpdf.extract_pdf_options import ExtractPDFOptions
from adobe.pdfservices.operation.pdfops.options.extractpdf.extract_element_type import ExtractElementType
from adobe.pdfservices.operation.execution_context import ExecutionContext
from adobe.pdfservices.operation.io.file_ref import FileRef
from adobe.pdfservices.operation.pdfops.extract_pdf_operation import ExtractPDFOperation
from nltk import WordNetLemmatizer, word_tokenize
from nltk.corpus import stopwords
from nltk.tag import pos_tag
# from transformers import (
#     TokenClassificationPipeline,
#     AutoModelForTokenClassification,
#     AutoTokenizer,
# )
from pptx import Presentation
from docx2pdf import convert
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Image
import os
import tempfile
from PIL import Image as PILImage
nlp = spacy.load("en_core_web_sm")

class CustomLogger:

    def set_process_id(self,process_id):
        self.process_id = process_id

    # Creates Log handlers
    def create_log_file(self,logger_name, log_file):
        try:
            logger = logging.getLogger(logger_name)
            logger.setLevel(logging.DEBUG)
            formatter = logging.Formatter('%(asctime)s %(levelname)s '+ self.process_id +' %(message)s')
            file_handler = logging.FileHandler(log_file)
            file_handler.setFormatter(formatter)
            logger.addHandler(file_handler)
            return logger
        except Exception as e:
            raise Exception(SupermatConstants.ERR_STRING_EXTRACT_ROLES)

    def log_closer(self,logger):
        try:
            handlers = logger.handlers[:]
            for handler in handlers:
                handler.close()
                logger.removeHandler(handler)
        except Exception as e:
            raise Exception(SupermatConstants.ERR_STR_LOG_CLOSE)

def log_create():
    try:
        pid = str(uuid.uuid4().hex)
        custom_process = CustomLogger()
        custom_process.set_process_id(pid)
        SUCCESS_LOG = custom_process.create_log_file(SUCCESS_NAME,SUCCESS_FILE)
        ERROR_LOG = custom_process.create_log_file(ERROR_NAME,ERROR_FILE)
        return SUCCESS_LOG, ERROR_LOG
    except Exception as e:
        raise Exception(SupermatConstants.ERR_STR_LOG_CREATE)

def convert_file_to_pdf(file_path):
    conversion_map = {
        '.pptx': convert_pptx_to_pdf,
        '.ppt': convert_pptx_to_pdf,
        '.doc': convert_docx_to_pdf,
        '.docx': convert_docx_to_pdf,
        '.odt': convert_docx_to_pdf  # Assuming the same function for odt as docx
    }
    
    _, file_extension = os.path.splitext(file_path.name)
    
    if file_extension in conversion_map:
        return conversion_map[file_extension](file_path)
    else:
        raise Exception('Unsupported file format')

def convert_docx_to_pdf(file_path):
    try:
        convert(file_path)
        return True
    except Exception as e:
        return False

def resize_image(image_path, max_width=500):
    with PILImage.open(image_path) as img:
        ratio = max_width / float(img.size[0])
        height = int(float(img.size[1]) * ratio)
        resized_img = img.resize((max_width, height))
        resized_img.save(image_path)

def convert_pptx_to_pdf(uploaded_file):
    try:
        temp_path = os.path.abspath(os.path.join(os.path.dirname(__name__)))
        pdf_file_path = os.path.join(str(temp_path), str(uploaded_file.name.split('.')[0]) + ".pdf")
        # Load the presentation
        prs = Presentation(uploaded_file.file)

        # Create a PDF
        pdf = SimpleDocTemplate(
            pdf_file_path,
            pagesize=letter
        )

        # Set the style
        styles = getSampleStyleSheet()
        styleN = styles['BodyText']

        # Collect slides' content
        content = []
        for slide_number, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Check if shape is an image
                    image = shape.image
                    image_bytes = image.blob
                    image_path = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                    image_path.write(image_bytes)
                    image_path.close()

                    # Resize the image
                    resize_image(image_path.name)

                    content.append(Image(image_path.name, width=400))  # Adjust width as needed
                elif hasattr(shape, "text"):
                    slide_content.append(shape.text)

            if slide_content:
                content.append(Paragraph("\n".join(slide_content), styleN))

        # Build the PDF
        pdf.build(content)
        return pdf_file_path
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error(f"Error converting {uploaded_file} to PDF: {str(e)} Trace Back: {str(trace_back)}")
        raise Exception(SupermatConstants.ERR_STRING_EXTRACT_ROLES)



def extract_roles(sentence):
    try:
        # Define the patterns for different roles
        # SUCCESS_LOG, ERROR_LOG = log_create()
        patterns = {
            "author": r'(author):',
            "operator": r'(operator):',
            "speaker": r'(speaker):',
            "moderator": r'(moderator):',
            "guest": r'(guest):'
        }
        roles = {}
        for role, pattern in patterns.items():
            match = re.search(pattern, sentence, re.IGNORECASE)
            if match:
                # Extract the role name after the colon
                roles[role] = sentence[match.end():].strip()
        return roles
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error("Error while Extracing the role: "+str(e)+". Trace Back: "+str(trace_back))
        raise Exception(SupermatConstants.ERR_STRING_EXTRACT_ROLES)

def extract_timestamp(text):
    try:
        # Extract timestamp from text
        SUCCESS_LOG, ERROR_LOG = log_create()
        pattern = rf'{SupermatConstants.TIME_STAMP_REGEX}'
        match = re.search(pattern, text)
        return match.group(1) if match else None
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error("Error while Extracing the timestamp" + str(e)+". Trace Back: "+ str(trace_back))
        raise Exception(SupermatConstants.ERR_STRING_EXTRACT_TIMESTAMP)

def extract_keywords_spacy(sentence):
    try:
        SUCCESS_LOG, ERROR_LOG = log_create()
        doc = nlp(sentence)
        # Extract meaningful words (nouns, adjectives, verbs, and adverbs)
        keywords = [token.lemma_ for token in doc if token.pos_ in ["NOUN", "ADJ", "VERB", "ADV"]]
        filtered_words = [word for word in keywords if len(word) >= 4]
        # Remove duplicates and return
        # SUCCESS_LOG.info("Spacy Extraction Successful")
        return list(set(filtered_words))
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error("Error while Extracing the keywords using Spacy: "+str(e)+". Trace Back: "+ str(trace_back))
        raise Exception(SupermatConstants.ERR_STRING_SPACY_EXTRACT)

def extract_meaningful_words(sentence):
    try:
        SUCCESS_LOG, ERROR_LOG = log_create()
        # Tokenize the sentence
        tokens = word_tokenize(sentence.lower())
        # Perform POS tagging
        tagged_tokens = pos_tag(tokens)
        # Extract nouns, verbs, adjectives, and adverbs
        meaningful_words = [word for word, tag in tagged_tokens if tag.startswith(('NN', 'VB', 'JJ', 'RB'))]
        filtered_words = [word for word in meaningful_words if len(word) >= 4]
        return list(set(filtered_words))
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error("Error while Extracing the keywords using NLTK: "+str(e)+". Trace Back: "+ str(trace_back))
        raise Exception(SupermatConstants.ERR_STRING_NLTK_EXTRACT)

def extract_keywords_nltk(sentence):
    try:
        SUCCESS_LOG, ERROR_LOG = log_create()
        # Tokenize the sentence
        tokens = word_tokenize(sentence.lower())
        # Remove stopwords and punctuation
        stop_words = set(stopwords.words('english'))
        tokens = [word for word in tokens if word.isalnum() and word not in stop_words]
        # Lemmatize words
        lemmatizer = WordNetLemmatizer()
        tokens = [lemmatizer.lemmatize(word) for word in tokens]
        # SUCCESS_LOG.info("NLTK extraction Successful")
        return tokens
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error("Error while Extracing the keywords using NLTK: "+str(e)+". Trace Back: "+ str(trace_back))
        raise Exception(SupermatConstants.ERR_STRING_NLTK_EXTRACT)

# Define keyphrase extraction pipeline
# class KeyphraseExtractionPipeline(TokenClassificationPipeline):
#     def __init__(self, model_name, *args, **kwargs):
#         super().__init__(
#             model=AutoModelForTokenClassification.from_pretrained(model_name),
#             tokenizer=AutoTokenizer.from_pretrained(model_name),
#             *args,
#             **kwargs
#         )

#     def postprocess(self, all_outputs):
#         results = super().postprocess(
#             all_outputs=all_outputs,
#             aggregation_strategy=AggregationStrategy.FIRST,
#         )
#         return np.unique([result.get("word").strip() for result in results])

# # Load the model outside the function to reuse it
# MODEL_NAME = SupermatConstants.HUGGING_FACE_MODEL_NAME
# EXTRACTOR = KeyphraseExtractionPipeline(MODEL_NAME)

# def hugging_face_extractor(sentence):
#     try:
#         cleaned_sentence = sentence.replace("\n", " ")
#         # Process the sentence
#         keywords = EXTRACTOR(cleaned_sentence)
#         return keywords
#     except Exception as e:
#         trace_back = traceback.format_exc()
#         ERROR_LOG.error(f"Error while extracting keywords using Hugging Face: {str(e)}. Trace Back: {str(trace_back)}")
#         raise Exception(SupermatConstants.ERR_STRING_HUGGING_FACE_EXTRACT)


def remove_similar_words(word_list, threshold=0.8):
    try:
        SUCCESS_LOG, ERROR_LOG = log_create()
        cleaned_list = []
        word_set = set(word_list)

        for word in word_set:
            if all(difflib.SequenceMatcher(None, word, w).ratio() < threshold for w in cleaned_list):
                cleaned_list.append(word)
        return list(set(cleaned_list))
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error(f"Error while removing the similar words from the keywords : {str(e)}. Trace Back: {str(trace_back)}")
        raise Exception(SupermatConstants.ERR_STR_REMOVE_SPECIAL_CHAR)

def extract_annotations(sentence):
    try:
        # Process the sentence
        doc = nlp(sentence)
        # Extract annotations (entities)
        annotations = [ent.text for ent in doc.ents]
        return list(set(annotations))
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error(f"Error extracting annotations: {str(e)}. Trace Back: {str(trace_back)}")
        raise Exception(SupermatConstants.ERR_STR_EXTRACt_ANNOT)

def remove_special_characters_from_list(lst):
    try:
        SUCCESS_LOG, ERROR_LOG = log_create()
        # Create translation table
        table = str.maketrans(string.punctuation, ' ' * len(string.punctuation))
        # Remove special characters from each string in the list
        cleaned_lst = [s.translate(table) for s in lst]
        # Remove extra spaces
        cleaned_lst = [' '.join(s.split()) for s in cleaned_lst]
        # Remove empty strings
        cleaned_lst = [s for s in cleaned_lst if s]
        # Remove similar words
        cleaned_lst = remove_similar_words(list(set(cleaned_lst)))
        return cleaned_lst
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error(f"Error while removing the special characters from the keywords: {str(e)}. Trace Back: {str(trace_back)}")
        raise Exception(SupermatConstants.ERR_STR_REMOVE_SPECIAL_CHAR)


def is_pdf(file_path):
    try:
        SUCCESS_LOG, ERROR_LOG = log_create()
        response = file_path.name.lower().endswith('.pdf')
        if file_path.content_type == 'application/pdf' and response:
            return True
        return False
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error("Error while checking the type: "+ str(e)+ ". Trace Back: "+str(trace_back))
        raise Exception(SupermatConstants.ERR_STRING_PDF_CHECK)

def parse_file(parsed_json, file_name, request_id):
    try:
        SUCCESS_LOG, ERROR_LOG = log_create()
        output = []
        existing_texts = set()
        def add_prefix_to_sentences(section_number_, passage_data, element_):
            keys_to_check = ['Bounds', 'Font', 'HasClip', 'Lang', 'ObjectID', 'Page', 'Path', 'TextSize', 'attributes']
            properties = {key: element_.get(key) for key in keys_to_check if key in element_}
            keywords = set(extract_keywords_spacy(passage_data)).union(extract_meaningful_words(passage_data))
            output.append({
                'type': 'Text',
                'structure': f'{section_number_}.{passage_number}.0',
                'text': passage_data,
                'key': remove_special_characters_from_list(keywords),
                'properties': properties,
                'sentences': [],
                'speaker': extract_roles(passage_data),
                'document': file_name,
                'timestamp': extract_timestamp(passage_data),
                'annotations': extract_annotations(passage_data)
            })
            sentences = re.split(rf'{SupermatConstants.SPLIT_SENTENCE_REGEX}', passage_data)
            sentences = [s.strip() for s in sentences if s]
            if len(sentences) > 1:
                for i, sentence in enumerate(sentences, start=1):
                    keywords = set(extract_keywords_spacy(sentence)).union(extract_meaningful_words(sentence))
                    output[-1]['sentences'].append({
                        'type': 'Text',
                        'structure': f'{section_number_}.{passage_number}.{i}',
                        'text': sentence,
                        'key': remove_special_characters_from_list(keywords),
                        'properties': properties
                    })
        section_number = 0
        passage_number = 0
        figure_count = 0
        for element in tqdm(parsed_json.get('elements',''), desc="Processing"):
            path_type = element.get('Path','')[11]
            if path_type == 'H':
                section_number += 1
                passage_number = 0
                add_prefix_to_sentences(section_number, element.get('Text',''), element)
            elif path_type == 'P':
                passage_number += 1
                text = element.get('Text','')
                if element.get('Path','').endswith('ParagraphSpan'):
                    pattern = re.compile(rf"{re.escape(element.get('Path',''))}\[\d+\]")
                    text += ''.join(elem.get('Text','') for elem in parsed_json.get('elements','') if pattern.match(elem.get('Path','')) and elem.get('Path','') not in existing_texts)
                add_prefix_to_sentences(section_number, text, element)
            elif 'Figure' in element.get('Path',''):
                passage_number += 1
                output.append({
                    'type': 'Image',
                    'structure': f'{section_number}.{passage_number}.0',
                    'figure': f'FIGURE {figure_count}',
                    'Bounds': element.get('Bounds',''),
                    'ObjectId': element.get('ObjectID',''),
                    'Page': element.get('Page',''),
                    'Path': element.get('Path',''),
                    'attributes': element.get('attributes','')
                })
                figure_count += 1
            elif path_type == 'L':
                match = re.match(r'^.*?/LI', element.get('Path',''))
                if match:
                    length_until_li = len(match.group())
                    text = element.get('Text','')
                    for elem in parsed_json.get('elements',''):
                        if elem.get('Path','')[11] == 'L' and elem.get('Path','')[:length_until_li] + "/Lbl" == element.get('Path',''):
                            section_number += 1
                            passage_number = 0
                            text += elem.get('Text','')
                    add_prefix_to_sentences(section_number, text, element)
            else:
                if 'Text' in element and element.get('Path','') not in existing_texts:
                    text = element.get('Text','')
                    pattern = re.compile(rf"{re.escape(element.get('Path',''))}\[\d+\]")
                    text += ''.join(elem.get('Text','') for elem in parsed_json.get('elements','') if pattern.match(elem.get('Path','')) and 'Text' in elem and elem.get('Path','') not in existing_texts)
                    passage_number += 1
                    add_prefix_to_sentences(section_number, text, element)
        SUCCESS_LOG.info(SupermatConstants.ERR_STR_PARSE_PDF.format(request_id=request_id))
        return output
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error(f"Error while parsing the file: parse_file {str(e)}. Traceback: {str(trace_back)}")
        raise Exception(SupermatConstants.ERR_STING_PDF_PARSER)

def adobe_pdf_parser(upload_file, request_id, is_original=False):
    try:
        SUCCESS_LOG, ERROR_LOG = log_create()
        if is_original:
            # Read the uploaded file
            uploaded_file_data = upload_file.file.read()
            file_name = upload_file.name
        else:
            with open(upload_file, 'rb') as file:
                uploaded_file_data = file.read()
                file_name = file.name
        
        # Setup Adobe credentials
        credentials = Credentials.service_principal_credentials_builder() \
            .with_client_id(SupermatConstants.ADOBE_CLIENT_ID) \
            .with_client_secret(SupermatConstants.ADOBE_CLIENT_SECRET) \
            .build()

        # Create an ExecutionContext using credentials and create a new operation instance
        execution_context = ExecutionContext.create(credentials)
        extract_pdf_operation = ExtractPDFOperation.create_new()

        # Set operation input from a source file
        source = FileRef.create_from_stream(BytesIO(uploaded_file_data), "application/pdf")
        extract_pdf_operation.set_input(source)

        # Build ExtractPDF options and set them into the operation
        extract_pdf_options = ExtractPDFOptions.builder() \
            .with_element_to_extract(ExtractElementType.TEXT) \
            .build()
        extract_pdf_operation.set_options(extract_pdf_options)

        # Execute the operation and get the result
        result = extract_pdf_operation.execute(execution_context)
        binary_stream = result.get_as_stream()

        # Extract JSON data from the zip file
        with zipfile.ZipFile(BytesIO(binary_stream), 'r') as zip_file:
            file_name = zip_file.namelist()[0]
            file_content = zip_file.read(file_name)
        
        json_data = json.loads(file_content)

        parsed_json = parse_file(json_data, file_name, request_id)

        if not is_original:
            try:
                if os.path.exists(file.name):
                    os.remove(file.name)
            except Exception as e:
                trace_back = traceback.format_exc()
                ERROR_LOG.critical(f"Error deleting {file.name}: {e}. Traceback: {str(trace_back)}")
        
        SUCCESS_LOG.info(SupermatConstants.ERR_STR_PARSE_PDF.format(request_id=request_id))
        return parsed_json
    except (ServiceApiException, ServiceUsageException, SdkException) as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.critical(f"Adobe Internal Exception: adobe_pdf_parser {str(e)}. Traceback: {str(trace_back)}")
        raise Exception(SupermatConstants.ERR_STRING_ADOBE_PARSER)
    except Exception as e:
        trace_back = traceback.format_exc()
        ERROR_LOG.error(f"Error while parsing the pdf: adobe_pdf_parser {str(e)}. Traceback: {str(trace_back)}")
        raise Exception(SupermatConstants.ERR_STRING_ADOBE_PARSER)